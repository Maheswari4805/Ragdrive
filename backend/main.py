import os
import uvicorn
from fastapi import FastAPI, HTTPException, Body
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional
import google.generativeai as genai
import fitz  # PyMuPDF
from pptx import Presentation
import io
import requests
import numpy as np

app = FastAPI(title="Ragdrive Python Backend")

# Enable CORS for the Expo app
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False, # Credentials cannot be True when Origins is wildcard "*"
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configuration from environment or keys
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    # Use a fallback key for testing if specified, otherwise error out
    raise ValueError("GEMINI_API_KEY Environment Variable not found on Render!")

genai.configure(api_key=GEMINI_API_KEY)

# --- Document Extraction Helpers ---

def extract_text(content: bytes, mime_type: str) -> str:
    if "pdf" in mime_type:
        doc = fitz.open(stream=content, filetype="pdf")
        return "".join([page.get_text() for page in doc])
    elif "presentation" in mime_type:
        prs = Presentation(io.BytesIO(content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        return content.decode("utf-8", errors="ignore")

# --- Manual RAG Core ---

class SimpleRag:
    def __init__(self):
        self.file_data = {} # file_id -> {"chunks": [], "embeddings": []}

    def compute_embedding(self, text: str):
        result = genai.embed_content(
            model="models/text-embedding-004",
            content=text,
            task_type="retrieval_document"
        )
        return result['embedding']

    def process_file(self, file_id: str, text: str):
        # Chunk text every 1000 characters
        chunks = [text[i:i+1000] for i in range(0, len(text), 800)]
        embeddings = []
        for chunk in chunks:
            try:
                embeddings.append(self.compute_embedding(chunk))
            except:
                continue
        self.file_data[file_id] = {"chunks": chunks, "embeddings": embeddings}

    def get_relevant_context(self, file_id: str, query: str):
        if file_id not in self.file_data:
            return ""
        
        query_emb = genai.embed_content(
            model="models/text-embedding-004",
            content=query,
            task_type="retrieval_query"
        )['embedding']
        
        data = self.file_data[file_id]
        dots = [np.dot(query_emb, chunk_emb) for chunk_emb in data["embeddings"]]
        top_indices = np.argsort(dots)[-5:][::-1]
        
        return "\n\n".join([data["chunks"][i] for i in top_indices])

rag = SimpleRag()

# --- API Endpoints ---

class ChatRequest(BaseModel):
    file_id: Optional[str] = None
    file_name: Optional[str] = "General Chat"
    mime_type: Optional[str] = "text/plain"
    access_token: Optional[str] = None
    query: str
    history: Optional[List[dict]] = []

@app.post("/chat")
async def chat_with_rag(request: ChatRequest):
    try:
        model_name = 'gemini-1.5-flash-latest'
        model = genai.GenerativeModel(model_name)
        
        # --- AGENTIC RAG STEP 1: Reasoning ---
        # First, we ask the brain if it needs a search for this specific query
        reasoning_prompt = f"""You are a RAG reasoning agent. 
        Analyze the USER_QUERY and decide if there is a document to search.
        
        USER_QUERY: "{request.query}"
        HAS_DOCUMENT: {'Yes' if request.file_id else 'No'}
        
        Respond with exactly one word: 'SEARCH' if we should scan the document, or 'CHAT' if this is a general greeting or non-document question.
        """
        reasoning_result = model.generate_content(reasoning_prompt).text.strip().upper()
        
        context = ""
        # --- AGENTIC RAG STEP 2: Tool Execution ---
        if 'SEARCH' in reasoning_result and request.file_id:
            # Automate file check & vector search
            if request.file_id not in rag.file_data:
                if not request.access_token:
                    raise HTTPException(status_code=400, detail="Access token required for file RAG")
                
                headers = {"Authorization": f"Bearer {request.access_token}"}
                url = f"https://www.googleapis.com/drive/v3/files/{request.file_id}?alt=media"
                response = requests.get(url, headers=headers)
                
                if response.status_code == 200:
                    text = extract_text(response.content, request.mime_type)
                    rag.process_file(request.file_id, text)
            
            context = rag.get_relevant_context(request.file_id, request.query)

        # --- AGENTIC RAG STEP 3: Final Response Generation ---
        if context:
            final_prompt = f"""You are the Drivedrop RAG Expert. 
            I have searched the document and found these relevant facts:
            {context}
            
            Answer the user's question accurately using these facts.
            USER_QUESTION: {request.query}
            """
        else:
            final_prompt = f"""You are the Drivedrop RAG Expert. 
            This is a general chat interaction without document search.
            
            USER_QUESTION: {request.query}
            """

        # Handle Chat History
        contents = []
        if request.history:
            for entry in request.history:
                role = "model" if entry.get("role") == "assistant" else "user"
                contents.append({"role": role, "parts": [entry.get("content", "")]})
        
        contents.append({"role": "user", "parts": [final_prompt]})
        
        result = model.generate_content(contents)
        return {"answer": result.text, "reasoning": reasoning_result}

    except Exception as e:
        print(f"Error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/")
def health_check():
    return {"status": "ok"}

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
